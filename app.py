import streamlit as st
import hashlib
import json
import os
import re
import openai
import pandas as pd
import pptx
import PyPDF2
from dotenv import load_dotenv
from pinecone import Pinecone, ServerlessSpec
import asyncio
from langchain.chains import RetrievalQA  
from langchain_openai import ChatOpenAI
from langchain_openai import OpenAIEmbeddings
from langchain_pinecone import PineconeVectorStore

# Load environment variables from .env file
load_dotenv()

# Read API keys from environment variables
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
pinecone_env = os.getenv("PINECONE_ENVIRONMENT")

# Initialize OpenAI
openai.api_key = OPENAI_API_KEY
client = openai

# Initialize Pinecone instance
pc = Pinecone(api_key=PINECONE_API_KEY)

# Function to hash passwords
def hash_password(password):
    return hashlib.sha256(password.encode()).hexdigest()

# Load persistent storage
def load_storage():
    if os.path.exists("storage.json"):
        with open("storage.json", "r") as f:
            return json.load(f)
    return {'users': {}, 'collections': {}, 'uploaded_files': {}}

# Save persistent storage
def save_storage():
    with open("storage.json", "w") as f:
        json.dump({
            'users': st.session_state.users,
            'collections': st.session_state.collections,
            'uploaded_files': st.session_state.uploaded_files
        }, f, indent=4)

# Initialize persistent storage
storage = load_storage()

# Initialize session state
if 'users' not in st.session_state:
    st.session_state.users = storage['users']
if 'current_user' not in st.session_state:
    st.session_state.current_user = None
if 'collections' not in st.session_state:
    st.session_state.collections = storage['collections']
if 'uploaded_files' not in st.session_state:
    st.session_state.uploaded_files = storage['uploaded_files']
if 'selected_collection' not in st.session_state:
    st.session_state.selected_collection = None
if 'page' not in st.session_state:
    st.session_state.page = 'login'

# Initialize Pinecone index state
if 'index' not in st.session_state:
    st.session_state.index = None

# User registration
def register(username, password):
    if username in st.session_state.users:
        return False
    st.session_state.users[username] = hash_password(password)
    st.session_state.collections[username] = []
    st.session_state.uploaded_files[username] = {}
    save_storage()
    return True

# User login
def login(username, password):
    if username in st.session_state.users and st.session_state.users[username] == hash_password(password):
        st.session_state.current_user = username
        return True
    return False

def add_collection(username, collection_name, index_name, uploaded_files):
    collections = st.session_state.collections.get(username, [])

    for collection in collections:
        if collection['name'] == collection_name:
            return "Collection name already exists."
        if collection['index_name'] == index_name:
            return "Index name already exists."

    result = asyncio.run(process_and_upsert_files(username, collection_name, index_name, uploaded_files))
    if result != "Success":
        return result

    collections.append({'name': collection_name, 'index_name': index_name})
    st.session_state.collections[username] = collections
    st.session_state.uploaded_files[username][collection_name] = [file.name for file in uploaded_files]

    st.write(f"Debug: Collection '{collection_name}' created with index '{index_name}'.")  # Log collection creation
    save_storage()

    return "Success"

async def process_and_upsert_files(username, collection_name, index_name, uploaded_files):
    try:
        if index_name not in pc.list_indexes().names():
            pc.create_index(
                name=index_name,
                dimension=1536,
                metric='cosine',
                spec=ServerlessSpec(cloud='aws', region='us-east-1')
            )

        # Wait for the index to be ready
        while not pc.describe_index(index_name).status['ready']:
            await asyncio.sleep(1)

        # Assign the index to session state
        st.session_state.index = pc.Index(index_name)
        st.write(f"Debug: Initialized Pinecone index: '{index_name}'")  # Log the index initialization

        for uploaded_file in uploaded_files:
            if uploaded_file.type == "text/plain":
                await process_txt_file(uploaded_file, collection_name)
            elif uploaded_file.type == "application/pdf":
                await process_pdf_file(uploaded_file, collection_name)
            elif uploaded_file.type == "text/csv":
                await process_csv_file(uploaded_file, collection_name)
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                await process_ppt_file(uploaded_file, collection_name)
            else:
                st.warning(f"Unsupported file type: {uploaded_file.type}")

        return "Success"
    except Exception as e:
        st.error(f"An error occurred: {e}")
        return f"Failed to create index or upsert files: {e}"

# Process and upsert a text file
async def process_txt_file(file, collection_name):
    text = file.read().decode("utf-8")
    await upsert_to_pinecone(text, collection_name)

# Process and upsert a PDF file
async def process_pdf_file(file, collection_name):
    reader = PyPDF2.PdfFileReader(file)
    text = ""
    for page_num in range(reader.getNumPages()):
        page = reader.getPage(page_num)
        try:
            text += page.extract_text()
        except AttributeError:
            text += page.extractText()
    await upsert_to_pinecone(text, collection_name)

# Process and upsert a CSV file
async def process_csv_file(file, collection_name):
    df = pd.read_csv(file)
    for _, row in df.iterrows():
        text = " ".join(row.astype(str).tolist())
        await upsert_to_pinecone(text, collection_name)

# Process and upsert a PPT file
async def process_ppt_file(file, collection_name):
    presentation = pptx.Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    await upsert_to_pinecone(text, collection_name)

# Upsert text into Pinecone
async def upsert_to_pinecone(text, collection_name):
    chunks = [text[i:i + 2000] for i in range(0, len(text), 2000)]
    for chunk in chunks:
        chunk = chunk.replace("\n", " ")
        try:
            response = client.embeddings.create(input=[chunk], model="text-embedding-ada-002")
            embedding = response.data[0].embedding
            
            if isinstance(embedding, list) and all(isinstance(x, float) for x in embedding):
                id = hashlib.sha256(chunk.encode()).hexdigest()
                st.session_state.index.upsert(vectors=[{"id": id,"values": embedding,"metadata": {"text": chunk}}])
            else:
                st.error("Failed to generate a valid embedding vector for upserting.")
        except Exception as e:
            st.error(f"An error occurred during embedding creation: {e}")

def query_pinecone_index_and_format_answer(question):
    try:
        st.write(f"Received question: '{question}'")  # Log the received question

        # Fetch the corresponding index name for the selected collection
        selected_collection = st.session_state.selected_collection
        index_name = None
        for collection in st.session_state.collections.get(st.session_state.current_user, []):
            if collection['name'] == selected_collection:
                index_name = collection['index_name']
                break
        
        if not index_name:
            st.error("Index name for the selected collection could not be found.")
            return None, None

        st.write(f"Using index name: '{index_name}'")  # Log the index name being used

        # Initialize the index if it hasn't been done already
        if 'index' not in st.session_state or st.session_state.index is None:
            st.session_state.index = pc.Index(index_name)  # Initialize index using the retrieved index name
            #st.write(f"Debug: Initialized Pinecone index: '{index_name}'")  # Log the index initialization

        # Proceed with the LLM and Pinecone-based retriever
        llm = ChatOpenAI(
            openai_api_key=OPENAI_API_KEY,
            model_name="gpt-3.5-turbo",
            temperature=0.0
        )

        knowledge = PineconeVectorStore.from_existing_index(
            index_name=index_name,
            embedding=OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY)
        )

        qa = RetrievalQA.from_chain_type(
            llm=llm,
            chain_type="stuff",
            retriever=knowledge.as_retriever()
        )

        # Use the RetrievalQA chain to answer the question
        #st.write("Retrieving answer...")  # Debugging statement
        answer = qa.run(question)

        if answer:
            source_text = "Relevant documents or information extracted."  # Replace this with actual source text if needed
            return answer, source_text
        else:
            st.warning("No relevant information found.")
            return None, None

    except Exception as e:
        st.error(f"An error occurred while querying Pinecone: {e}")
        return None, None

# UI components and navigation
def show_login_page():
    st.title("Login")
    username = st.text_input("Username")
    password = st.text_input("Password", type="password")
    if st.button("Login"):
        if username and password:
            if login(username, password):
                st.success(f"Welcome, {username}!")
                st.session_state.page = 'create_collection'
            else:
                st.error("Invalid username or password.")
        else:
            st.error("Username and Password cannot be empty.")

    st.write("Don't have an account?")
    if st.button("Register"):
        st.session_state.page = 'register'

def show_registration_page():
    st.title("Register")
    username = st.text_input("Username")
    password = st.text_input("Password", type="password")
    if st.button("Register"):
        if username and password:
            if register(username, password):
                st.success("Registration successful! You can now log in.")
                st.session_state.page = 'login'
            else:
               
                                st.error("Username already exists.")
        else:
            st.error("Username and Password cannot be empty.")
    
    if st.button("Back to Login"):
        st.session_state.page = 'login'

def show_create_collection_page():
    st.title("Create Collection")
    collection_name = st.text_input("Collection Name")
    index_name = st.text_input("Index Name")
    uploaded_files = st.file_uploader("Upload Files", accept_multiple_files=True)

    if st.button("Create Collection"):
        if collection_name and index_name and uploaded_files:
            result = add_collection(st.session_state.current_user, collection_name, index_name, uploaded_files)
            if result == "Success":
                st.success(f"Collection '{collection_name}' created successfully!")
            else:
                st.error(result)
        else:
            st.error("All fields are required.")

    st.write("Or select an existing collection:")
    collections = st.session_state.collections.get(st.session_state.current_user, [])
    if collections:
        collection_names = [collection['name'] for collection in collections]
        selected_collection = st.selectbox("Select Collection", collection_names)
        if st.button("Select"):
            st.session_state.selected_collection = selected_collection
            st.session_state.page = 'ask_question'

import asyncio
import streamlit as st

async def upsert_embedding(blended_sentence, collection_name):
    await upsert_to_pinecone(blended_sentence, collection_name)
    st.success("Embedding created and upserted successfully!")  # Success message

def show_ask_question_page():
    st.title(f"Ask a Question in '{st.session_state.selected_collection}'")
    question = st.text_area("Your Question")

    if st.button("Ask"):
        if question:
            answer, source_text = query_pinecone_index_and_format_answer(question)
            if answer:
                st.write("### Answer:")
                st.write(answer)
                st.write("### Source Documents:")
                st.write(source_text)

                # Ask for user satisfaction with a selection box
                satisfied = st.selectbox("Are you satisfied with the answer?", ["Select None", "Yes", "No"])

                # Checkbox for user to decide whether to upsert
                if satisfied == "Yes":
                    upsert_choice = st.checkbox("Do you want to upsert the blended sentence?", value=False)

                    if upsert_choice and st.button("Upsert Blended Sentence"):
                        blended_sentence = f"The user asked: '{question}' and the answer provided was: '{answer}'."
                        st.write("Generating embedding for the blended sentence...")

                        # Run the feedback handling asynchronously
                        result = asyncio.run(handle_feedback(blended_sentence, st.session_state.selected_collection))
                        st.success(result)

                elif satisfied == "No":
                    st.write("You indicated that you are not satisfied. No embedding will be created.")

                elif satisfied == "Select None":
                    st.write("No action will be taken.")
            else:
                st.error("No relevant information found.")
        else:
            st.error("Please enter a question.")

    # Logout Button
    if st.button("Logout"):
        st.session_state.current_user = None
        st.session_state.page = 'login'
        st.success("Logged out successfully")


  
# Note: Ensure to call this function in your main navigation method
def navigation():
    st.title("PROJECT MANAGEMENT AI TOOL")
    if st.session_state.page == 'login':
        show_login_page()
    elif st.session_state.page == 'register':
        show_registration_page()
    elif st.session_state.page == 'create_collection':
        show_create_collection_page()
    elif st.session_state.page == 'ask_question':
        show_ask_question_page()

navigation()
